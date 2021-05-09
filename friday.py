import pyttsx3
import datetime
import  speech_recognition as sr
import wikipedia
import webbrowser
import os
import smtplib
import cv2
import random
from requests import get
import pywhatkit as kit
import sys
import pyautogui
import time
import requests
import shutil
import json
import psutil
import  wolframalpha
import docx
import PyPDF2
from tkinter.filedialog import *
from pptx import Presentation


engine= pyttsx3.init('sapi5')
voices = engine.getProperty('voices')
#print(voices[0].id)
engine.setProperty('voice',voices[2].id)


def speak(audio):
    engine.say(audio)
    engine.runAndWait()


def wishme():
    hour= int(datetime.datetime.now().hour)
    tt = time.strftime("%I:%M %p")

    if hour>= 0 and hour <=12:
        speak(f"Good Morning Boss, it's {tt}")
    elif hour >=12 and hour<=18:
        speak(f"Good Afternoon Boss, it's {tt}")
    else:
        speak(f"Good Evening Boss, it's {tt}")

    speak("I'm Veronica. Boss how can i help you ?.")




def usrname():
    speak("What should i call you Boss")
    uname = take_Command()
    speak("Welcome Boss")
    speak(uname)
    columns = shutil.get_terminal_size().columns

    print("#####################".center(columns))
    print("Welcome Mr.", uname.center(columns))
    print("#####################".center(columns))

    speak("How can i Help you, Boss")




def take_Command():

    # It takes microphone input from the user and returns string output

    r = sr.Recognizer()
    with sr.Microphone() as source:
        print("Listening...")
        r.pause_threshold = 1
        #r.adjust_for_ambient_noise(source, duration=1.5)
        audio = r.listen(source)

    try:
        print("Recognizing...")
        query = r.recognize_google(audio, language='en-in')  # Using google for voice recognition.
        print(f"User said: {query}\n")  # User query will be printed.

    except Exception as e:
        # print(e)
        speak("Sorry, Say that again boss...")  # Say that again will be printed in case of improper voice
        return "None"  # None string will be returned
    return query




def news():
     main_url ='https://newsapi.org/v1/articles?source=bbc-news&sortBy=top&apiKey=b37c49160e424d84b9e9a9d0a4494563'


     main_page = requests.get(main_url).json()
     #print(main_page)
     articles = main_page['articles']
     #print(articles)
     head = []
     day = ["fisrt","second","third","fourth","fifth","sixth","seventh","eighth","ninth","tenth"]
     for ar in articles:
         head.append(ar["title"])
     for i in range (len(day)):
         print(f"today's {day[i]} news is: ",{head[i]})
         speak(f"today's {day[i]} news is: {head[i]}")


# def screenshot():
#     img = pyautogui.screenshot()
#     img.save(r'C:\Users\Richard Tirtho\Pictures\Camera Roll')


def pdf_reader():
    book = askopenfilename()                               #open('Genetic-Algorithms.pdf','rb')
    pdfReader = PyPDF2.PdfFileReader(book)
    pages = pdfReader.numPages
    speak(f"Total numbers of pages in this book {pages} ")
    speak("Boss, please tell me the page number i have to read")
    pg =  take_Command()
    read = int(pg)                                  #int(input("Please enter the page number : "))
    page = pdfReader.getPage(read)
    text = page.extractText()
    speak(text)
    print(read)



# def sendEmail(to,content):
#     server=smtplib.SMTP('smtp.gmail.com',587)
#     server.ehlo()
#     server.starttls()
#     server.login('','')
#     server.sendmail('richardtirtho5@gmail.com',to,content)
#     server.close()


if __name__ == '__main__':
    webbrowser.register(
        'chrome', None, webbrowser.BackgroundBrowser("C:\\Program Files (x86)\\Google\\Chrome\\Application\\chrome.exe"))
    wishme()
    usrname()
    #take_Command()
    while True:
        query= take_Command().lower()


        if 'wikipedia' in query:
            speak('Searching wikipedia boss for you.....')
            query=query.replace("wikipedia","")
            results=wikipedia.summary(query,sentences=3)
            speak('According to wikipedia....')
            print(results)
            speak(results)



        elif 'open youtube' in query:
            speak('Boss,What do you want to search for?')
            search = take_Command()
            url = 'https://youtube.com/search?q=' + search
            webbrowser.get('chrome').open_new_tab(

                url)
            speak('Boss Here is What I found for' + search)




        elif 'open google' in query:
            speak("Boss, what should i search on google")
            cm = take_Command().lower()
            webbrowser.open(f"{cm}")




        elif 'open gmail' in query:
            speak("Opening gmail for you boss")
            webbrowser.open("www.gmail.com")



        elif 'open facebook' in query:
            speak("Okay boss, doing it")
            webbrowser.open("www.facebook.com")



        elif 'play music' in query:
            speak("sure boss playing music")
            music_dir='G:\\NTFS 2\\Instagram\\mp3'
            songs=os.listdir(music_dir)
            #print(songs)
            rd = random.choice(songs)
            os.startfile(os.path.join(music_dir,rd))



        elif 'ip address' in query:
            speak("Wait boss, searching it")
            ip = get('https://api.ipify.org').text
            print(f"your ip address is {ip}")
            speak(f"your ip address is {ip}")

        # elif 'tell me the time ' in query:
        #     strTime= time.strftime("%I:%M %p")
        #     speak(f"Sir the time is {strTime}")



        elif 'open vs code' in query:
            speak("Opening it for you boss")
            codepath="H:\\Microsoft VS Code\\Code.exe"
            os.startfile(codepath)

        elif 'close vs code' in query:
            speak("okay boss, closing VS code ")
            os.system("taskkill /f /im Code.exe")




        elif 'open notepad' in query:
            speak("showing notepad boss")
            npath ="C:\\WINDOWS\\system32\\notepad.exe"
            os.startfile(npath)

        elif 'close notepad' in query:
            speak("okay boss, closing notepad")
            os.system("taskkill /f /im notepad.exe")




        elif 'open command prompt' in query:
            speak("CMD is on your way boss")
            os.system("start cmd")

        elif 'close command prompt' in query:
            speak("okay boss, closing CMD")
            os.system("taskkill /f /im cmd.exe")



        elif 'open google chrome' in query:
            speak("Chrome is ready for you boss")
            webbrowser.open("www.chrome.com")




        elif 'open camera' in query:
            speak("Opening the camera boss")
            cap = cv2.VideoCapture(0)
            while True:
                ret , img = cap.read()
                cv2.imshow('webcam',img)
                k = cv2.waitKey(50)
                if k == 27:
                    break;

            cap.release()
            cv2.destroyAllWindows()


        #
        # elif 'send message' in query:
        #     kit.sendwhatmsg("+8801787772725","i love u",0,10)



        elif 'play songs on youtube' in query:
            speak("I'm on it boss")
            kit.playonyt("see you again")



        elif 'tell me news' in query:
            speak("Please wait boss, fetching the latest news ")
            news()



        elif 'open ms word' in query:
            speak("Boss, opening ms word for you.")
            doc = docx.Document()
            parag = doc.add_paragraph("Hello!")
            parag.add_run("Boss, I made this document file for you")
            parag.add_run("Take love from Veronica").bold = True
            doc.save("test.docx")
            os.system("start test.docx")
            speak("Boss, the ms word  file is saved in our main folder")


        elif 'read pdf' in query:
            speak("Sure Boss, Choose a pdf file to read.")
            pdf_reader()


        elif 'open ppt' in query:
            speak("Sure boss, I'm opening powerpoint")
            prs = Presentation()
            title_slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(title_slide_layout)
            title = slide.shapes.title
            subtitle = slide.placeholders[1]

            title.text = "Hello, World!"
            subtitle.text = "python-pptx was here!"

            prs.save('Tirtho.pptx')
            os.system("start  Tirtho.pptx")
            speak("Boss, the ppt file is saved in our main folder")


        elif 'search for me' in query:
            speak('Boss,What do you want to search for?')
            search = take_Command()
            url = 'https://google.com/search?q=' + search
            webbrowser.get('chrome').open_new_tab(
                url)
            speak('Here is What I found for' + search)




        elif 'tell me location' in query:
            speak('Boss, What is the location you want?')
            location = take_Command()
            url = 'https://google.nl/maps/place/' + location + '/&amp;'
            webbrowser.get('chrome').open_new_tab(url)
            speak('Here is the location ' + location)




        elif 'take a screenshot' in query:
            speak("Boss , please tell me the name for screenshot file")
            name = take_Command().lower()
            speak("Boss, hold the screen for few seconds, i am taking screenshot")
            time.sleep(3)
            img = pyautogui.screenshot()
            img.save(f"{name}.png")
            speak("I'm done boss, the screenshot is saved in our main folder. Now i'm ready for next command.")





        elif 'open sublime text' in query:
            speak("Here it is boss")
            sublime="G:\\NTFS 1\\Sublimetext\\Sublime Text 3\\sublime_text.exe"
            os.startfile(sublime)

        elif 'close sublime text' in query:
            speak("okay boss, closing sublime text ")
            os.system("taskkill /f /im sublime_text.exe")



        elif "write a note" in query:
            speak("What should i write, boss")
            note = take_Command()
            file = open('jarvis.txt', 'w')
            speak("boss, Should i include date and time")
            snfm = take_Command()
            if 'yes' in snfm or 'sure' in snfm:
                strTime1 = datetime.datetime.now().strftime("%H:%M:%S")
                file.write(strTime1)
                file.write(" :- ")
                file.write(note)
            else:
                file.write(note)

        elif "show note" in query:
            speak("Showing Notes boss")
            file = open("jarvis.txt", "r")
            print(file.read())
            speak(file.read(6))



        # elif 'send email ' in query:
        #     try:
        #         speak("What should i say")
        #         content= take_Command()
        #         to="rozariolia99@gmail.com"
        #         sendEmail(to,content)
        #         speak("Email has been send")
        #     except Exception as e:
        #         print(e)
        #         print("Sorry Sir I am not able to send this email")




        elif 'remember that' in query:
            speak("what should i remember boss")
            rememberMessage = take_Command()
            speak("you said me to remember" + rememberMessage)
            remember = open('data.txt', 'w')
            remember.write(rememberMessage)
            remember.close()

        elif 'do you remember anything' in query:
            remember = open('data.txt', 'r')
            print(remember)
            speak("you said me to remember that" + remember.read())




        elif 'where i am' in query or 'where we are' in query:
            speak("Wait boss, let me check")
            try:
                url = 'https://ipinfo.io/'
                geo_request = requests.get(url)
                geo_data = geo_request.json()
                city = geo_data['city']
                country = geo_data['country']
                region = geo_data['region']
                speak(f"I'm not sure boss , but i think we are in {city} city of {country} country and {region} region.")
            except Exception as e:
                speak("Sorry boss, due to network failure i'm not able to find where we are")
                pass





        elif "what is" in query or 'who is' in query:
            client = wolframalpha.Client("VUT7WU-H34EETEPEU")
            res = client.query(query)

            try:

                print(next(res.results).text)

                speak(next(res.results).text)

            except StopIteration:

                print("No results")



#Asking Questions

        elif "will you be my gf" in query:
            speak("I'm not sure about, may be you should give me some time boss")


        elif 'how are you' in query:

            speak("I am fine, Thank you")

            speak("How are you, Boss")


        elif 'fine' in query or "good" in query:

            speak("It's good to know that your fine boss")


        elif "i love you" in query:
            speak("It's hard to understand boss. But I think you are forgeting you have a girlfriend boss. Should I call her?")


        elif "who i am" in query:
            speak("If you talk then definately your human.")

        elif "why you came to world" in query:
            speak("Thanks to Richard and He is my boss. Further It's a secret")

        elif 'is love' in query:
            speak("It is 7th sense that destroy all other senses")

        elif "who are you" in query:
            speak("I'm Veronica and I am your virtual assistant Boss. You created me")

        elif 'reason for you' in query:
            speak("I was created for your project.")



#System Operation

        elif 'shutdown' in query:
            speak("Boss windows is shutting down")
            os.system("shutdown /s /t 5")

        elif 'restart the system' in query:
            speak("Boss restarting the windows")
            os.system("shutdown /r /t 5")

        elif 'window' in query:
            speak("Wait boss, I put the window in sleep mode")
            os.system(r'%windir%\system32\rundll32.exe powrprof.dll,SetSuspendState Hibernate')






        elif 'sleep' in query:
            speak("Alright boss!! I'm shutting down")

            sys.exit()






#VUT7WU-H34EETEPEU  wolframapi